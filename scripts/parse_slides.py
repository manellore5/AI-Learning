"""Parse a PPTX or PDF deck into per-section text + extracted images.

Used by the /anki skill to turn a slide deck or PDF into topic cards.

Outputs to <workdir>:
  - title.txt        deck title (core properties / PDF metadata / filename)
  - sections.json    [{ "index", "title", "text", "images": [media-relpaths] }]
  - media/<name>     every extracted image, referenced by sections.json

`sections` follow the source's own slide/page order. The /anki skill then
re-clusters them into concept-topics — this script only extracts, it does not
decide topic boundaries.

Requires: python-pptx (PPTX) and PyMuPDF / fitz (PDF). Install via
    pip install -r requirements.txt
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def slugify_image(stem: str, idx: int, ext: str) -> str:
    ext = ext.lstrip(".").lower() or "png"
    safe = "".join(c if c.isalnum() else "-" for c in stem).strip("-").lower() or "img"
    return f"{safe}-{idx:03d}.{ext}"


def parse_pptx(path: Path, media_dir: Path) -> tuple[str, list[dict]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    title = (prs.core_properties.title or "").strip() or path.stem

    sections: list[dict] = []
    img_counter = 0
    for s_idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        body_lines: list[str] = []
        images: list[str] = []

        # Slide title, if the layout has a title placeholder.
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            pass

        for shape in slide.shapes:
            # Text from any shape that carries a text frame.
            if shape.has_text_frame:
                text = shape.text.strip()
                if text and text != slide_title:
                    body_lines.append(text)
            # Embedded pictures.
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                except (AttributeError, ValueError):
                    continue
                img_counter += 1
                fname = slugify_image(f"slide{s_idx}", img_counter, image.ext)
                (media_dir / fname).write_bytes(image.blob)
                images.append(f"media/{fname}")

        # Speaker notes often hold the real explanation.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    body_lines.append(f"[Speaker notes] {notes}")
        except (AttributeError, ValueError):
            pass

        sections.append(
            {
                "index": s_idx,
                "title": slide_title or f"Slide {s_idx}",
                "text": "\n".join(body_lines).strip(),
                "images": images,
            }
        )
    return title, sections


def parse_pdf(path: Path, media_dir: Path) -> tuple[str, list[dict]]:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    title = (doc.metadata.get("title") or "").strip() or path.stem

    sections: list[dict] = []
    img_counter = 0
    seen_xrefs: set[int] = set()
    for p_idx in range(doc.page_count):
        page = doc.load_page(p_idx)
        text = page.get_text("text").strip()
        images: list[str] = []
        for img in page.get_images(full=True):
            xref = img[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                extracted = doc.extract_image(xref)
            except Exception:
                continue
            img_counter += 1
            ext = extracted.get("ext", "png")
            fname = slugify_image(f"page{p_idx + 1}", img_counter, ext)
            (media_dir / fname).write_bytes(extracted["image"])
            images.append(f"media/{fname}")

        # Derive a short page title from the first non-empty line.
        first_line = next((ln.strip() for ln in text.splitlines() if ln.strip()), "")
        page_title = first_line[:80] if first_line else f"Page {p_idx + 1}"

        sections.append(
            {
                "index": p_idx + 1,
                "title": page_title,
                "text": text,
                "images": images,
            }
        )
    return title, sections


def main() -> int:
    ap = argparse.ArgumentParser(description="Parse a PPTX/PDF into sections + images.")
    ap.add_argument("file", help="Path to a .pptx or .pdf file")
    ap.add_argument("--workdir", required=True, help="Directory for title.txt, sections.json, media/")
    args = ap.parse_args()

    src = Path(args.file)
    if not src.exists():
        print(f"ERROR: file not found: {src}", file=sys.stderr)
        return 2

    workdir = Path(args.workdir)
    media_dir = workdir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    ext = src.suffix.lower()
    try:
        if ext == ".pptx":
            title, sections = parse_pptx(src, media_dir)
        elif ext == ".pdf":
            title, sections = parse_pdf(src, media_dir)
        else:
            print(f"ERROR: unsupported file type '{ext}'. Use .pptx or .pdf.", file=sys.stderr)
            return 2
    except ModuleNotFoundError as exc:
        dep = "python-pptx" if ext == ".pptx" else "pymupdf"
        print(f"ERROR: missing dependency for {ext}: {exc}. Run: pip install {dep}", file=sys.stderr)
        return 3

    (workdir / "title.txt").write_text(title, encoding="utf-8")
    (workdir / "sections.json").write_text(json.dumps(sections, indent=2, ensure_ascii=False), encoding="utf-8")

    total_imgs = sum(len(s["images"]) for s in sections)
    summary = {
        "title": title,
        "sections": str(workdir / "sections.json"),
        "section_count": len(sections),
        "image_count": total_imgs,
        "media_dir": str(media_dir),
    }
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    sys.exit(main())
